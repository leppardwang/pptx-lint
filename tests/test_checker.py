"""Basic tests for pptx-lint rules."""

from pptx import Presentation
from pptx.util import Inches

from pptx_lint.utils import emu_to_inches
from pptx_lint import rules


def _make_dummy_shapes_info():
    """Return a shapes_info list with two non-overlapping text boxes."""
    return [
        {
            'idx': 0,
            'bounds': (0.5, 0.5, 4.0, 0.5, 4.5, 1.0),
            'texts': ['Left box'],
            'font_info': [{'text': 'Left box', 'size': 14}],
            'type': 'TEXT_BOX',
        },
        {
            'idx': 1,
            'bounds': (5.0, 0.5, 4.0, 0.5, 9.0, 1.0),
            'texts': ['Right box'],
            'font_info': [{'text': 'Right box', 'size': 14}],
            'type': 'TEXT_BOX',
        },
    ]


# ── overflow ──────────────────────────────────────────────────

def test_overflow_no_error():
    infos = _make_dummy_shapes_info()
    issues = rules.overflow.check(10.0, 7.5, infos)
    assert len(issues) == 0


def test_overflow_right_edge():
    infos = _make_dummy_shapes_info()
    infos[0]['bounds'] = (0.5, 0.5, 10.0, 0.5, 10.5, 1.0)  # 10.5 > 10.0
    issues = rules.overflow.check(10.0, 7.5, infos)
    assert len(issues) == 1
    assert 'right' in issues[0]['issues'][0]


def test_overflow_bottom_edge():
    infos = _make_dummy_shapes_info()
    infos[0]['bounds'] = (0.5, 7.6, 4.0, 0.5, 4.5, 8.1)  # 8.1 > 7.5
    issues = rules.overflow.check(10.0, 7.5, infos)
    assert len(issues) == 1
    assert 'bottom' in issues[0]['issues'][0]


# ── overlap ───────────────────────────────────────────────────

def test_overlap_none():
    infos = _make_dummy_shapes_info()
    issues = rules.overlap.check(infos)
    assert len(issues) == 0


def test_overlap_detected():
    # Two boxes at the exact same position
    infos = _make_dummy_shapes_info()
    infos.append({
        'idx': 2,
        'bounds': (0.5, 0.5, 4.0, 0.5, 4.5, 1.0),
        'texts': ['Overlapping'],
        'font_info': [],
        'type': 'TEXT_BOX',
    })
    issues = rules.overlap.check(infos)
    assert len(issues) >= 1


# ── fake table ────────────────────────────────────────────────

def test_fake_table_none_for_list():
    """A vertical list (same column, many rows) is NOT a fake table."""
    infos = []
    for i in range(6):
        infos.append({
            'idx': i,
            'bounds': (0.5, 1.0 + i * 0.5, 4.0, 0.4, 4.5, 1.4 + i * 0.5),
            'texts': [f'Item {i}'],
            'font_info': [],
            'type': 'TEXT_BOX',
        })
    issues = rules.fake_table.check(infos)
    assert len(issues) == 0


def test_fake_table_detected_for_grid():
    """A 2×2 grid IS a fake table."""
    infos = []
    for row in range(2):
        for col in range(2):
            infos.append({
                'idx': row * 2 + col,
                'bounds': (
                    0.5 + col * 4.5,
                    1.0 + row * 0.6,
                    4.0,
                    0.5,
                    4.5 + col * 4.5,
                    1.5 + row * 0.6,
                ),
                'texts': [f'R{row}C{col}'],
                'font_info': [],
                'type': 'TEXT_BOX',
            })
    issues = rules.fake_table.check(infos)
    assert len(issues) == 1


# ── font size ─────────────────────────────────────────────────

def test_small_font_detected():
    infos = _make_dummy_shapes_info()
    infos[0]['font_info'] = [{'text': 'tiny', 'size': 6}]  # 6pt < 8pt
    issues = rules.font_size.check(infos)
    assert len(issues) == 1
    assert issues[0]['font_size'] == 6


# ── end-to-end via checker ────────────────────────────────────

def _create_minimal_pptx(tmp_path):
    """Create a minimal .pptx for integration testing."""
    from pptx import Presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5),
                                      Inches(9.0), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Hello pptx-lint"
    p.font.size = Pt(14)
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return str(path)


def test_end_to_end(tmp_path):
    path = _create_minimal_pptx(tmp_path)
    from pptx_lint.checker import check_pptx
    results, name = check_pptx(path)
    assert len(results) == 1
    assert 'Hello' in results[0]['text_content'][0]
    # No errors expected for this well-formed slide
    assert len(results[0]['errors']) == 0
