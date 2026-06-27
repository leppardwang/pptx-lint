"""Core checker: orchestrate all lint rules against one or more PPTX files."""

import os
from pptx import Presentation

from .utils import emu_to_inches, emu_to_pt
from . import rules


def get_shape_bounds(shape):
    """Return (left, top, width, height, right, bottom) in inches."""
    try:
        left = emu_to_inches(shape.left)
        top = emu_to_inches(shape.top)
        w = emu_to_inches(shape.width)
        h = emu_to_inches(shape.height)
        return left, top, w, h, left + w, top + h
    except Exception:
        return None


def get_shape_texts(shape):
    """Extract plain text from a shape (text-frame or table)."""
    texts = []
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            t = p.text.strip()
            if t:
                texts.append(t)
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            parts = []
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    parts.append(t)
            if parts:
                texts.append(' | '.join(parts))
    return texts


def collect_shapes_info(slide):
    """Return list of shape info dicts for a single slide."""
    infos = []
    for idx, shape in enumerate(slide.shapes):
        bounds = get_shape_bounds(shape)
        if bounds is None:
            continue

        texts = get_shape_texts(shape)

        font_info = []
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    sz = p.font.size
                    font_info.append({
                        'text': p.text.strip(),
                        'size': emu_to_pt(sz) if sz else None,
                    })

        infos.append({
            'idx': idx,
            'bounds': bounds,
            'texts': texts,
            'font_info': font_info,
            'type': str(shape.shape_type),
        })
    return infos


def check_slide(slide, slide_num, sw, sh):
    """Run all lint rules against one slide.

    Returns a dict with keys:
        slide_num, errors, warnings, infos, text_content, is_empty
    """
    result = {
        'slide_num': slide_num,
        'total_shapes': len(slide.shapes),
        'errors': [],
        'warnings': [],
        'infos': [],
        'text_content': [],
        'is_empty': True,
    }

    shapes_info = collect_shapes_info(slide)
    for info in shapes_info:
        result['text_content'].extend(info['texts'])
    result['is_empty'] = len(result['text_content']) == 0

    # ── overflow (error) ──
    for issue in rules.overflow.check(sw, sh, shapes_info):
        result['errors'].append({
            'type': 'Overflow',
            'detail': (
                f"Shape#{issue['shape_idx']} ({issue['shape_name']}): "
                f"{' | '.join(issue['issues'])}"
            ),
            'bounds': (
                f"({issue['bounds'][0]:.2f}, {issue['bounds'][1]:.2f}, "
                f"{issue['bounds'][2]:.2f}, {issue['bounds'][3]:.2f})"
            ),
        })

    # ── empty slide (error) ──
    if result['total_shapes'] == 0:
        result['errors'].append({
            'type': 'EmptySlide',
            'detail': 'Slide has no shapes at all.',
        })

    # ── overlap (warning) ──
    for issue in rules.overlap.check(shapes_info):
        result['warnings'].append({
            'type': 'Overlap',
            'detail': (
                f"Shape#{issue['shapes'][0]} ({issue['names'][0]}) "
                f"overlaps Shape#{issue['shapes'][1]} ({issue['names'][1]}) "
                f"— overlap area: {issue['overlap_area']}"
            ),
        })

    # ── fake table (warning) ──
    for issue in rules.fake_table.check(shapes_info):
        samples = ', '.join(issue['sample_texts'])
        result['warnings'].append({
            'type': 'FakeTable',
            'detail': (
                f"{issue['shape_count']} text boxes form a grid-like layout "
                f"(e.g.: {samples}) — consider using ``add_table()`` instead."
            ),
        })

    # ── small font (info) ──
    for issue in rules.font_size.check(shapes_info):
        result['infos'].append({
            'type': 'SmallFont',
            'detail': (
                f"Shape#{issue['shape_idx']}: '{issue['text']}' "
                f"font size = {issue['font_size']:.0f}pt "
                f"(< {rules.font_size.MIN_FONT_SIZE}pt)"
            ),
        })

    return result


def check_pptx(filepath):
    """Run all checks on a single .pptx file.

    Returns list of slide-level result dicts.
    """
    prs = Presentation(filepath)
    sw = emu_to_inches(prs.slide_width)
    sh = emu_to_inches(prs.slide_height)

    all_results = []
    for si, slide in enumerate(prs.slides):
        all_results.append(check_slide(slide, si + 1, sw, sh))
    return all_results, os.path.basename(filepath)
